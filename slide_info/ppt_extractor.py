from pptx import Presentation
import os

def extract_shapes(slide):
    shapes = []
    for shape in slide.shapes:
        if shape.shape_type == 6:  # GroupShape
            shapes.extend(extract_group_shapes(shape))
        else:
            shapes.append(shape)
    return shapes

def extract_group_shapes(group_shape):
    shapes = []
    for shape in group_shape.shapes:
        if shape.shape_type == 6:  # Nested GroupShape
            shapes.extend(extract_group_shapes(shape))
        else:
            shapes.append(shape)
    return shapes

def extract_slide_data(ppt_file_path, config):
    presentation = Presentation(ppt_file_path)
    file_name = os.path.basename(ppt_file_path)
    
    slide_data = {
        'file_name': [],
        'slide_no': [],
        'title': [],
        'rfp': [],
        'page_no': [],
        'navigation': []
    }
    
    for slide_number, slide in enumerate(presentation.slides, start=0):
        if slide_number != 0:
            slide_data['file_name'].append(file_name)
            slide_data['slide_no'].append(slide_number)
            
            title_text = ''
            rfp_text = ''
            page_no_text = ''
            navigation_text = ''
            
            all_shapes = extract_shapes(slide)

            for shape in all_shapes:
                if shape.has_text_frame:
                    shape_name = shape.name
                    if shape_name == config['Title']:
                        title_text = file_name[0] + "." + shape.text
                    elif shape_name == config['Rfp']:
                        rfp_text = shape.text
                    elif shape_name == config['PageNo']:
                        page_no_text = shape.text
                    elif shape_name == config['Navigation']:
                        navigation_text = shape.text
        
            if page_no_text == '':
                page_no_text = create_page_no(slide_number, file_name)

            slide_data['title'].append(title_text)
            slide_data['rfp'].append(rfp_text)
            slide_data['page_no'].append(page_no_text)
            slide_data['navigation'].append(navigation_text)
  
    return slide_data

def create_page_no(slide_no, file_name):
    base_name = file_name[0]
    return f"{base_name}-{slide_no}"
