from pptx import Presentation
import pandas as pd
from pathlib import Path
import configparser
import os

# 설정 파일 읽기
config = configparser.ConfigParser()
config.read('./slide_info/config.ini')

conf_page_no = config['DEFAULT']['PageNo']
conf_title = config['DEFAULT']['Title']
conf_rfp = config['DEFAULT']['Rfp']
conf_eval_item = config['DEFAULT']['EvalItem']
conf_directory_path = config['DEFAULT']['DirectoryPath']

def extract_shapes(slide):
    shapes = []
    for shape in slide.shapes:
        if shape.shape_type == 6:  # GroupShape
            shapes.extend(extract_group_shapes(shape))
        else:
            shapes.append(shape)
    return shapes

def extract_group_shapes(group_shape):
    shapes = []
    for shape in group_shape.shapes:
        if shape.shape_type == 6:  # Nested GroupShape
            shapes.extend(extract_group_shapes(shape))
        else:
            shapes.append(shape)
    return shapes

def extract_slide_data(ppt_file_path):
    # 프레젠테이션 로드
    presentation = Presentation(ppt_file_path)
    
    # 파일명 추출
    file_name = os.path.basename(ppt_file_path)
    
    # 슬라이드 데이터 초기화
    slide_data = {
        'file_name': [],
        'slide_no': [],
        'title': [],
        'rfp': [],
        'page_no': [],
        'eval_item': []
    }
    
    # 슬라이드 내용 추출
    for slide_number, slide in enumerate(presentation.slides, start=0):
        if slide_number != 0:
            slide_data['file_name'].append(file_name)
            slide_data['slide_no'].append(slide_number)
            
            title_text = ''
            rfp_text = ''
            page_no_text = ''
            eval_item_text = ''
            
            all_shapes = extract_shapes(slide)

            for shape in all_shapes:
                if shape.has_text_frame:
                    shape_name = shape.name
                    if shape_name == conf_title:
                        title_text = file_name[0]+"."+shape.text
                    elif shape_name == conf_rfp:
                        rfp_text = shape.text
                    elif shape_name == conf_page_no:
                        page_no_text = shape.text
                    elif shape_name == conf_eval_item:
                        eval_item_text = shape.text
        
            if page_no_text == '':
                page_no_text = create_page_no(slide_number,file_name)

            slide_data['title'].append(title_text)
            slide_data['rfp'].append(rfp_text)
            slide_data['page_no'].append(page_no_text)
            slide_data['eval_item'].append(eval_item_text)
  
    return slide_data

def create_page_no(slide_no, file_name):
    base_name = file_name[0]
    return f"{base_name}-{slide_no}"


def save_to_excel(data, excel_file_path):
    df = pd.DataFrame(data)
    df.to_excel(excel_file_path, index=False)
    print(f'DataFrame이 {excel_file_path} 파일로 저장되었습니다.')

def list_ppt_files(directory_path):
    ppt_files = [str(file) for file in Path(directory_path).rglob('*.pptx') if file.is_file()]
    return sorted(ppt_files)

def merge_data(ppts_data):
    merge_data = {
        '파일명': [],
        '슬라이드번호': [],
        '제목': [],
        '평가항목': [],
        '페이지': [],
        'rfp': []
    }
    
    for data in ppts_data:
        merge_data['파일명'].extend(data['file_name'])
        merge_data['슬라이드번호'].extend(data['slide_no'])
        merge_data['제목'].extend(data['title'])
        merge_data['평가항목'].extend(data['eval_item'])
        merge_data['페이지'].extend(data['page_no'])
        merge_data['rfp'].extend(data['rfp'])
    
    return merge_data

def main():
    directory_path = conf_directory_path
    excel_file_path = os.path.join(directory_path, '제안서_슬라이드_정보.xlsx')
    
    ppt_files = list_ppt_files(directory_path)
    ppts_data = [extract_slide_data(ppt_file) for ppt_file in ppt_files]
    
    merged_data = merge_data(ppts_data)
    save_to_excel(merged_data, excel_file_path)

if __name__ == '__main__':
    main()
