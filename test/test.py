import os
import configparser
from pathlib import Path
from pptx import Presentation
import pandas as pd

# 설정 파일 읽기
config = configparser.ConfigParser()
config.read('config.ini')

conf_page_no = config['DEFAULT']['PageNo'].lower()
conf_title = config['DEFAULT']['Title'].lower()
conf_rfp = config['DEFAULT']['Rfp'].lower()
conf_eval_item = config['DEFAULT']['EvalItem'].lower()
conf_directory_path = config['DEFAULT']['DirectoryPath'].lower()

def create_page_no(slide_no, file_name):
    base_name = file_name.rsplit('.', 1)[0]
    return f"{base_name}-{slide_no}"

def extract_slide_texts(slide, config_mappings):
    texts = {key: '' for key in config_mappings.values()}
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape_name = shape.name.lower()
            if shape_name in config_mappings:
                texts[config_mappings[shape_name]] = shape.text
    return texts

def extract_slide_data(ppt_file_path, config_mappings):
    presentation = Presentation(ppt_file_path)
    file_name = os.path.basename(ppt_file_path)
    slide_data = {
        'file_name': [],
        'slide_no': [],
        'title': [],
        'rfp': [],
        'page_no': [],
        'eval_item': []
    }
    
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_data['file_name'].append(file_name)
        slide_data['slide_no'].append(slide_number)
        
        texts = extract_slide_texts(slide, config_mappings)
        
        if not texts['page_no']:
            texts['page_no'] = create_page_no(slide_number, file_name)
        
        slide_data['title'].append(texts['title'])
        slide_data['rfp'].append(texts['rfp'])
        slide_data['page_no'].append(texts['page_no'])
        slide_data['eval_item'].append(texts['eval_item'])
    
    return slide_data

def save_to_excel(data, excel_file_path):
    df = pd.DataFrame(data)
    df.to_excel(excel_file_path, index=False)
    print(f'DataFrame이 {excel_file_path} 파일로 저장되었습니다.')

def list_ppt_files(directory_path):
    return [str(file) for file in Path(directory_path).rglob('*.pptx') if file.is_file()]

def merge_data(ppts_data):
    merged_data = {
        '파일명': [],
        '슬라이드번호': [],
        '제목': [],
        '평가항목': [],
        '페이지': [],
        'rfp': []
    }
    
    for data in ppts_data:
        for key in merged_data.keys():
            merged_data[key].extend(data[key])
    
    return merged_data

def main():
    directory_path = conf_directory_path
    excel_file_path = os.path.join(directory_path, '제안서_슬라이드_정보.xlsx')
    
    config_mappings = {
        conf_page_no: 'page_no',
        conf_title: 'title',
        conf_rfp: 'rfp',
        conf_eval_item: 'eval_item'
    }
    
    ppt_files = list_ppt_files(directory_path)
    ppts_data = [extract_slide_data(ppt_file, config_mappings) for ppt_file in ppt_files]
    
    merged_data = merge_data(ppts_data)
    save_to_excel(merged_data, excel_file_path)

if __name__ == '__main__':
    main()
